import streamlit as st  
from pptx import Presentation  
from pptx.util import Inches, Emu  
from pptx.enum.text import PP_ALIGN  
from io import BytesIO  
  
# Constants  
PIC_START_SLIDE = 7  # Slide number where pictures start  
SPACE_HEIGHT = Inches(6.3)  # Space for images height  
SPACE_WIDTH = Inches(13)  # Space for images width  
TITLE_HEIGHT = Inches(0.52)  
TITLE_LEFT = Inches(0.2)  
TITLE_TOP = Inches(0.5)  # Adjust this value as necessary for vertical alignment  
TITLE_WIDTH = Inches(12)  # Adjust this value as necessary for width  
FIGURE_MARGIN = Inches(0.2)  # Margin between figures  
SLIDE_HEIGHT = Inches(7.5)  # Standard slide height  
SLIDE_WIDTH = Inches(13.33)  # Standard slide width  
  
def modify_text_box(shape):  
    shape.text = ''  
    shape.width = Inches(11.5)  
    shape.height = Inches(3.12)  
    shape.left = Inches(0.91)  
    shape.top = Inches(1.87)  
    if shape.has_text_frame:  
        tf = shape.text_frame  
        if tf.text:  # Ensure the text frame has text  
            try:  
                tf.fit_text(font_family=u'Verdana', max_size=60, bold=True)  
            except TypeError:  
                st.write("Error fitting text for shape ID:", shape.shape_id)  
  
def modify_title(shape):  
    shape.height = TITLE_HEIGHT  
    shape.width = TITLE_WIDTH  
    shape.left = TITLE_LEFT  
    shape.top = TITLE_TOP  
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  
    if shape.has_text_frame:  
        tf = shape.text_frame  
        if tf.text:  # Ensure the text frame has text  
            try:  
                tf.fit_text(font_family=u'Verdana', max_size=24, bold=True)  
            except TypeError:  
                st.write("Error fitting text for shape ID:", shape.shape_id)  
    return shape.height.inches  
  
def modify_picture(shape, top_margin, available_height):  
    shape.left = Inches(0.2)  
    shape.top = Inches(top_margin)  
    pic_ratio = shape.width / shape.height  
    space_ratio = SPACE_WIDTH / available_height  
  
    # Adjust length and width to fit slide  
    if pic_ratio > space_ratio:  
        shape.width = Emu(SPACE_WIDTH)  
        shape.height = Emu(SPACE_WIDTH / pic_ratio)  
    else:  
        shape.height = Emu(available_height)  
        shape.width = Emu(available_height * pic_ratio)  
  
    # Ensure the top position does not exceed the slide height  
    if shape.top + shape.height > SLIDE_HEIGHT:  
        shape.top = SLIDE_HEIGHT - shape.height  
  
    # Return the bottom position of the current shape  
    return shape.top + shape.height.inches + FIGURE_MARGIN  
  
def layout_pictures(slide, top_margin):  
    shapes = [shape for shape in slide.shapes if shape.shape_type == 13]  # Filter picture shapes  
    num_shapes = len(shapes)  
      
    if num_shapes == 0:  
        return top_margin  
  
    available_height = SPACE_HEIGHT - top_margin  
    max_height_per_shape = available_height / num_shapes - FIGURE_MARGIN  
  
    for shape in shapes:  
        top_margin = modify_picture(shape, top_margin, max_height_per_shape)  
        available_height = SPACE_HEIGHT - top_margin  # Recalculate available height for next figure  
  
    return top_margin  
  
# Streamlit UI for file upload  
st.title('PowerPoint Slide Modifier')  
uploaded_file = st.file_uploader("Choose a PPTX file", type="pptx")  
  
if uploaded_file is not None:  
    # Open pptx  
    prs = Presentation(uploaded_file)  
    slides = prs.slides  
  
    # Scan each slide  
    for slide in slides:  
        slide_index = slides.index(slide) + 1  
        slide_title = slide.shapes.title.text if slide.shapes.title else "No title"  
        st.write(f'\nslide number {slide_index}: {slide_title}')  
  
        top_margin = 0  
        for shape in slide.shapes:  
            # Show all objects in the slide  
            st.write(f'id: {shape.shape_id}, height: {round(shape.height.inches, 2)}, width: {round(shape.width.inches, 2)}, left: {round(shape.left.inches, 2)}', shape.shape_type, shape.name)  
  
            # Modify slide sections to section header-like format  
            if shape.shape_type == 17 and hasattr(shape, "text") and shape.text == '<change layout for title slide>':  
                modify_text_box(shape)  
  
            # Modify slide titles starting page 2  
            if shape.shape_type == 14 and slide_index != 1 and hasattr(shape, "text") and shape.text == slide.shapes.title.text:  
                top_margin = modify_title(shape)  
  
        # Modify pages starting pic_start for picture  
        if slide_index >= PIC_START_SLIDE and slide_title != 'Legends':  
            top_margin = layout_pictures(slide, top_margin)  
  
    # Save the modified presentation to an in-memory file  
    output = BytesIO()  
    prs.save(output)  
    output.seek(0)  
  
    # Provide a download link  
    st.download_button(  
        label="Download Modified PPTX",  
        data=output,  
        file_name="modified_presentation.pptx",  
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"  
    )  
